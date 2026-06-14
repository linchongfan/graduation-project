from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_arrow(slide, start_x, start_y, end_x, end_y, color=RGBColor(0, 120, 215), width=2):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector

def add_circle(slide, left, top, size, fill_color=RGBColor(0, 120, 215)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(2)
    return shape

def add_text_in_shape(slide, left, top, width, height, text, font_size=14, bold=False, color=RGBColor(0, 0, 0), fill_color=RGBColor(0, 120, 215)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape

def create_2opt_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    add_textbox(slide, Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6),
                "2-opt局部搜索算法演示", font_size=28, bold=True,
                color=RGBColor(0, 100, 180), alignment=PP_ALIGN.CENTER)
    
    # 城市位置（左右两侧相同）
    city_size = Inches(0.35)
    offset_x = Inches(0.175)
    offset_y = Inches(0.175)
    
    # 5个城市的位置 - 梯形分布，容易产生交叉
    city_positions = [
        (Inches(1.2), Inches(1.8)),   # A - 左上
        (Inches(2.8), Inches(1.5)),   # B - 中上
        (Inches(4.2), Inches(2.2)),   # C - 右上
        (Inches(2.2), Inches(3.0)),   # D - 左下
        (Inches(3.8), Inches(3.2)),   # E - 右下
    ]
    
    city_labels = ["A", "B", "C", "D", "E"]
    
    # ===== 左侧：原始路径（有交叉） =====
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(4), Inches(0.4),
                "原始路径（存在交叉）", font_size=16, bold=True, color=RGBColor(200, 50, 50))
    
    # 绘制城市
    for i, (x, y) in enumerate(city_positions):
        add_circle(slide, x, y, city_size, RGBColor(0, 120, 215))
        add_textbox(slide, x, y + Inches(0.05), city_size, Inches(0.25),
                    city_labels[i], font_size=12, bold=True, color=RGBColor(255, 255, 255),
                    alignment=PP_ALIGN.CENTER)
    
    # 原始连接顺序: A-B-C-D-E-A（会有交叉）
    original_connections = [(0, 1), (1, 2), (2, 3), (3, 4), (4, 0)]
    for i, j in original_connections:
        x1, y1 = city_positions[i][0] + offset_x, city_positions[i][1] + offset_y
        x2, y2 = city_positions[j][0] + offset_x, city_positions[j][1] + offset_y
        add_arrow(slide, x1, y1, x2, y2, RGBColor(200, 50, 50), 2)
    
    # 标注交叉点
    add_textbox(slide, Inches(2.8), Inches(2.3), Inches(1.5), Inches(0.35),
                "交叉!", font_size=14, bold=True, color=RGBColor(200, 50, 50))
    
    # 路径长度
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(4), Inches(0.35),
                "路径: A→B→C→D→E→A", font_size=13, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(0.5), Inches(3.9), Inches(4), Inches(0.35),
                "总距离: 12.5 单位", font_size=13, bold=True, color=RGBColor(200, 50, 50))
    
    # ===== 中间：箭头 =====
    add_textbox(slide, Inches(4.6), Inches(2.0), Inches(0.8), Inches(0.6),
                "→", font_size=40, bold=True, color=RGBColor(0, 150, 0),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.5), Inches(2.6), Inches(1), Inches(0.35),
                "2-opt", font_size=12, bold=True, color=RGBColor(0, 150, 0),
                alignment=PP_ALIGN.CENTER)
    
    # ===== 右侧：优化后路径 =====
    add_textbox(slide, Inches(5.5), Inches(1.0), Inches(4), Inches(0.4),
                "2-opt优化后（消除交叉）", font_size=16, bold=True, color=RGBColor(0, 150, 0))
    
    # 右侧城市位置（与左侧相同，但x偏移5单位）
    right_offset = Inches(5)
    for i, (x, y) in enumerate(city_positions):
        rx = x + right_offset
        add_circle(slide, rx, y, city_size, RGBColor(0, 180, 0))
        add_textbox(slide, rx, y + Inches(0.05), city_size, Inches(0.25),
                    city_labels[i], font_size=12, bold=True, color=RGBColor(255, 255, 255),
                    alignment=PP_ALIGN.CENTER)
    
    # 优化后连接顺序: A-B-E-D-C-A（反转C-D-E片段，消除交叉）
    optimized_connections = [(0, 1), (1, 4), (4, 3), (3, 2), (2, 0)]
    for i, j in optimized_connections:
        x1, y1 = city_positions[i][0] + right_offset + offset_x, city_positions[i][1] + offset_y
        x2, y2 = city_positions[j][0] + right_offset + offset_x, city_positions[j][1] + offset_y
        add_arrow(slide, x1, y1, x2, y2, RGBColor(0, 150, 0), 2)
    
    # 路径长度
    add_textbox(slide, Inches(5.5), Inches(3.6), Inches(4), Inches(0.35),
                "路径: A→B→E→D→C→A", font_size=13, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(5.5), Inches(3.9), Inches(4), Inches(0.35),
                "总距离: 10.8 单位", font_size=13, bold=True, color=RGBColor(0, 150, 0))
    
    # ===== 底部：算法步骤 =====
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.4),
                "2-opt操作步骤：", font_size=16, bold=True, color=RGBColor(0, 100, 180))
    
    # 步骤框
    steps = [
        ("① 选择两条边\n(B,C)和(D,E)", Inches(0.5), RGBColor(200, 120, 0)),
        ("② 计算反转后\n路径长度变化", Inches(2.7), RGBColor(0, 120, 215)),
        ("③ ΔL<0，执行\n反转操作", Inches(4.9), RGBColor(0, 150, 0)),
        ("④ 重复直到\n无法改进", Inches(7.1), RGBColor(150, 0, 150))
    ]
    
    for text, x, color in steps:
        add_text_in_shape(slide, x, Inches(5.0), Inches(1.8), Inches(0.8),
                         text, font_size=11, bold=True, color=RGBColor(255, 255, 255), fill_color=color)
    
    # 箭头连接步骤
    for i in range(3):
        x1 = steps[i][1] + Inches(1.8)
        x2 = steps[i + 1][1]
        add_arrow(slide, x1, Inches(5.4), x2, Inches(5.4), RGBColor(100, 100, 100), 1.5)
    
    # 反转操作说明
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(9), Inches(0.35),
                "核心：反转 [C→D→E] 为 [E→D→C]，消除边交叉，路径更短",
                font_size=12, bold=True, color=RGBColor(80, 80, 80))

def create_ga_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8),
                "遗传算法增强蚁群算法演示", font_size=32, bold=True,
                color=RGBColor(0, 100, 180), alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(2.5), Inches(0.5),
                "1. 锦标赛选择", font_size=18, bold=True, color=RGBColor(200, 120, 0))
    
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(2.5), Inches(0.3),
                "父代路径 A: 1-2-3-4-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(2.5), Inches(0.3),
                "父代路径 B: 1-3-2-5-4-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(2.5), Inches(0.3),
                "父代路径 C: 1-2-5-3-4-6", font_size=12, color=RGBColor(80, 80, 80))
    
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(2.5), Inches(0.3),
                "选择优质路径: A, C", font_size=12, bold=True, color=RGBColor(0, 150, 0))
    
    add_textbox(slide, Inches(3.2), Inches(1.2), Inches(3), Inches(0.5),
                "2. 顺序交叉(OX)", font_size=18, bold=True, color=RGBColor(0, 120, 215))
    
    add_textbox(slide, Inches(3.2), Inches(1.7), Inches(3), Inches(0.3),
                "父代A: 1-2-|3-4-|-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(3.2), Inches(2.0), Inches(3), Inches(0.3),
                "父代C: 1-2-|5-3-|-4-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(3.2), Inches(2.3), Inches(3), Inches(0.3),
                "子代:   1-2-|3-4-|-6-5", font_size=12, color=RGBColor(0, 150, 0))
    
    add_textbox(slide, Inches(3.2), Inches(2.8), Inches(3), Inches(0.3),
                "保留优秀基因，生成新路径", font_size=11, color=RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(3), Inches(0.5),
                "3. 变异操作", font_size=18, bold=True, color=RGBColor(150, 0, 150))
    
    add_textbox(slide, Inches(6.5), Inches(1.7), Inches(3), Inches(0.3),
                "变异前: 1-2-3-4-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(6.5), Inches(2.0), Inches(3), Inches(0.3),
                "变异后: 1-2-5-4-3-6", font_size=12, color=RGBColor(150, 0, 150))
    add_textbox(slide, Inches(6.5), Inches(2.3), Inches(3), Inches(0.3),
                "随机交换两个城市", font_size=11, color=RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(6.5), Inches(2.8), Inches(3), Inches(0.3),
                "前10%精英释放2倍信息素", font_size=11, bold=True, color=RGBColor(200, 50, 50))
    
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.5),
                "遗传算法与蚁群算法融合流程：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    flow_steps = [
        ("蚁群构路", Inches(1), RGBColor(0, 120, 215)),
        ("筛选优质路径", Inches(2.8), RGBColor(200, 120, 0)),
        ("交叉操作", Inches(4.6), RGBColor(0, 150, 0)),
        ("变异操作", Inches(6.4), RGBColor(150, 0, 150)),
        ("精英强化", Inches(8.2), RGBColor(200, 50, 50))
    ]
    
    for text, x, color in flow_steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(1.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    for i in range(len(flow_steps) - 1):
        x1 = flow_steps[i][1] + Inches(1.6)
        x2 = flow_steps[i + 1][1]
        add_arrow(slide, x1, Inches(4.5), x2, Inches(4.5), RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
                "遗传算法增强效果：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    effects = [
        "• 最优距离降低 3.95%（单独使用）",
        "• 与2-opt组合降低 8.13%（最优组合）",
        "• 维持种群多样性，增强全局搜索能力",
        "• 精英策略加速收敛，避免退化"
    ]
    
    for i, effect in enumerate(effects):
        add_textbox(slide, Inches(0.5), Inches(5.7 + i * 0.3), Inches(9), Inches(0.3),
                    effect, font_size=12, color=RGBColor(80, 80, 80))

def main():
    prs = Presentation(r'C:\Users\帆\Desktop\毕业设计\林重帆答辩PPT.pptx')
    
    create_2opt_slide(prs)
    create_ga_demo(prs)
    
    output_path = r'C:\Users\帆\Desktop\毕业设计\林重帆答辩PPT_算法演示v2.pptx'
    prs.save(output_path)
    print(f"已保存到: {output_path}")
    print(f"总幻灯片数: {len(prs.slides)}")

if __name__ == "__main__":
    main()
