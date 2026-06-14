from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_arrow(slide, start_x, start_y, end_x, end_y, color=RGBColor(0, 120, 215)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def add_circle(slide, left, top, size, fill_color=RGBColor(0, 120, 215)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0, 120, 215)
    return shape

def create_2opt_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8),
                "2-opt局部搜索算法演示", font_size=32, bold=True,
                color=RGBColor(0, 100, 180), alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.5),
                "原始路径（存在交叉）", font_size=18, bold=True, color=RGBColor(200, 50, 50))
    
    cities = [
        (Inches(1), Inches(2), "A"),
        (Inches(2.5), Inches(1.8), "B"),
        (Inches(4), Inches(2.5), "C"),
        (Inches(2), Inches(3.2), "D"),
        (Inches(3.5), Inches(3), "E")
    ]
    
    for x, y, label in cities:
        add_circle(slide, x, y, Inches(0.4), RGBColor(0, 120, 215))
        add_textbox(slide, x, y + Inches(0.05), Inches(0.4), Inches(0.3),
                    label, font_size=14, bold=True, color=RGBColor(255, 255, 255),
                    alignment=PP_ALIGN.CENTER)
    
    connections = [(0, 1), (1, 2), (2, 3), (3, 4), (4, 0)]
    for i, j in connections:
        x1, y1 = cities[i][0] + Inches(0.2), cities[i][1] + Inches(0.2)
        x2, y2 = cities[j][0] + Inches(0.2), cities[j][1] + Inches(0.2)
        add_arrow(slide, x1, y1, x2, y2, RGBColor(200, 50, 50))
    
    add_textbox(slide, Inches(3), Inches(2.1), Inches(1.5), Inches(0.4),
                "交叉点", font_size=14, bold=True, color=RGBColor(200, 50, 50))
    
    add_textbox(slide, Inches(4.5), Inches(2.2), Inches(1), Inches(0.5),
                "→", font_size=36, bold=True, color=RGBColor(0, 150, 0),
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(5.5), Inches(1.2), Inches(4), Inches(0.5),
                "2-opt优化后", font_size=18, bold=True, color=RGBColor(0, 150, 0))
    
    cities_opt = [
        (Inches(6), Inches(2), "A"),
        (Inches(7.5), Inches(1.8), "B"),
        (Inches(9), Inches(2.5), "C"),
        (Inches(7.5), Inches(3.2), "D"),
        (Inches(6.5), Inches(3), "E")
    ]
    
    for x, y, label in cities_opt:
        add_circle(slide, x, y, Inches(0.4), RGBColor(0, 180, 0))
        add_textbox(slide, x, y + Inches(0.05), Inches(0.4), Inches(0.3),
                    label, font_size=14, bold=True, color=RGBColor(255, 255, 255),
                    alignment=PP_ALIGN.CENTER)
    
    connections_opt = [(0, 1), (1, 3), (3, 4), (4, 2), (2, 0)]
    for i, j in connections_opt:
        x1, y1 = cities_opt[i][0] + Inches(0.2), cities_opt[i][1] + Inches(0.2)
        x2, y2 = cities_opt[j][0] + Inches(0.2), cities_opt[j][1] + Inches(0.2)
        add_arrow(slide, x1, y1, x2, y2, RGBColor(0, 150, 0))
    
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.5),
                "2-opt核心步骤：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    steps = [
        "1. 选择两条边 (i,i+1) 和 (j,j+1)",
        "2. 计算反转片段后的路径长度变化 ΔL",
        "3. 若 ΔL < 0，执行反转操作",
        "4. 重复直到无法改进"
    ]
    
    for i, step in enumerate(steps):
        add_textbox(slide, Inches(0.5), Inches(4.7 + i * 0.35), Inches(9), Inches(0.35),
                    step, font_size=14, color=RGBColor(80, 80, 80))
    
    add_textbox(slide, Inches(5.5), Inches(4.2), Inches(4), Inches(0.5),
                "优化效果：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    effects = [
        "• 最优距离降低 7.79%",
        "• 收敛次数从 8.3 降至 2.0",
        "• 时间复杂度 O(n²)"
    ]
    
    for i, effect in enumerate(effects):
        add_textbox(slide, Inches(5.5), Inches(4.7 + i * 0.35), Inches(4), Inches(0.35),
                    effect, font_size=14, color=RGBColor(80, 80, 80))

def create_ga_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8),
                "遗传算法增强蚁群算法演示", font_size=32, bold=True,
                color=RGBColor(0, 100, 180), alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(2.5), Inches(0.5),
                "1. 锦标赛选择", font_size=18, bold=True, color=RGBColor(200, 120, 0))
    
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(2.5), Inches(0.3),
                "父代路径 A: 1-2-3-4-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(2.5), Inches(0.3),
                "父代路径 B: 1-3-2-5-4-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(2.5), Inches(0.3),
                "父代路径 C: 1-2-5-3-4-6", font_size=12, color=RGBColor(80, 80, 80))
    
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(2.5), Inches(0.3),
                "选择优质路径: A, C", font_size=12, bold=True, color=RGBColor(0, 150, 0))
    
    add_textbox(slide, Inches(3.2), Inches(1.2), Inches(3), Inches(0.5),
                "2. 顺序交叉(OX)", font_size=18, bold=True, color=RGBColor(0, 120, 215))
    
    add_textbox(slide, Inches(3.2), Inches(1.7), Inches(3), Inches(0.3),
                "父代A: 1-2-|3-4-|-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(3.2), Inches(2.0), Inches(3), Inches(0.3),
                "父代C: 1-2-|5-3-|-4-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(3.2), Inches(2.3), Inches(3), Inches(0.3),
                "子代:   1-2-|3-4-|-6-5", font_size=12, color=RGBColor(0, 150, 0))
    
    add_textbox(slide, Inches(3.2), Inches(2.8), Inches(3), Inches(0.3),
                "保留优秀基因，生成新路径", font_size=11, color=RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(3), Inches(0.5),
                "3. 变异操作", font_size=18, bold=True, color=RGBColor(150, 0, 150))
    
    add_textbox(slide, Inches(6.5), Inches(1.7), Inches(3), Inches(0.3),
                "变异前: 1-2-3-4-5-6", font_size=12, color=RGBColor(80, 80, 80))
    add_textbox(slide, Inches(6.5), Inches(2.0), Inches(3), Inches(0.3),
                "变异后: 1-2-5-4-3-6", font_size=12, color=RGBColor(150, 0, 150))
    add_textbox(slide, Inches(6.5), Inches(2.3), Inches(3), Inches(0.3),
                "随机交换两个城市", font_size=11, color=RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(6.5), Inches(2.8), Inches(3), Inches(0.3),
                "前10%精英释放2倍信息素", font_size=11, bold=True, color=RGBColor(200, 50, 50))
    
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.5),
                "遗传算法与蚁群算法融合流程：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    flow_steps = [
        ("蚁群构路", Inches(1), RGBColor(0, 120, 215)),
        ("筛选优质路径", Inches(2.8), RGBColor(200, 120, 0)),
        ("交叉操作", Inches(4.6), RGBColor(0, 150, 0)),
        ("变异操作", Inches(6.4), RGBColor(150, 0, 150)),
        ("精英强化", Inches(8.2), RGBColor(200, 50, 50))
    ]
    
    for text, x, color in flow_steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(1.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    for i in range(len(flow_steps) - 1):
        x1 = flow_steps[i][1] + Inches(1.6)
        x2 = flow_steps[i + 1][1]
        add_arrow(slide, x1, Inches(4.5), x2, Inches(4.5), RGBColor(100, 100, 100))
    
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
                "遗传算法增强效果：", font_size=18, bold=True, color=RGBColor(0, 100, 180))
    
    effects = [
        "• 最优距离降低 3.95%（单独使用）",
        "• 与2-opt组合降低 8.13%（最优组合）",
        "• 维持种群多样性，增强全局搜索能力",
        "• 精英策略加速收敛，避免退化"
    ]
    
    for i, effect in enumerate(effects):
        add_textbox(slide, Inches(0.5), Inches(5.7 + i * 0.3), Inches(9), Inches(0.3),
                    effect, font_size=12, color=RGBColor(80, 80, 80))

def main():
    prs = Presentation(r'C:\Users\帆\Desktop\毕业设计\林重帆答辩PPT.pptx')
    
    create_2opt_demo(prs)
    create_ga_demo(prs)
    
    output_path = r'C:\Users\帆\Desktop\毕业设计\林重帆答辩PPT_算法演示版.pptx'
    prs.save(output_path)
    print(f"已保存到: {output_path}")
    print(f"总幻灯片数: {len(prs.slides)}")

if __name__ == "__main__":
    main()
